"""PPTX 解析器模块。

PPT / PPTX 在企业知识库里也非常常见：

- 产品宣讲材料
- 培训课件
- 方案评审稿
- 架构分享文档

这类文档的难点在于：

- 内容被切散在多个 slide
- 标题和正文的层级很重要
- 表达往往是短句和 bullet points，而不是完整段落

因此这里的策略是：

1. 逐页读取 slide
2. 优先提取标题
3. 再按 shape 顺序提取文本
4. 输出为带 slide 标记和标题的结构化文本
"""

from __future__ import annotations

from pathlib import Path

from core.ingestion.cleaners.text_cleaner import clean_text
from core.models.document import Document

from .base import BaseParser


def _safe_slide_title(slide) -> str:
    """尽量从 slide 中提取标题。

    PPT 的标题通常是整页语义最强的锚点：

    - 后续切块时可以作为页级标题
    - 检索命中后能帮助用户快速理解片段来自哪一页
    """

    title_shape = getattr(slide.shapes, "title", None)
    if title_shape is not None:
        text = clean_text(getattr(title_shape, "text", "")).strip()
        if text:
            return text
    return ""


def _shape_lines(shape) -> list[str]:
    """从 shape 中尽量提取带层级的文本行。

    相比普通文档段落，PPT 更像“标题 + 若干短句 bullet”。
    这里保留 bullet 缩进层级，是为了尽量维持课件/汇报材料原本的结构感。
    """

    text_frame = getattr(shape, "text_frame", None)
    if text_frame is None:
        text = clean_text(getattr(shape, "text", "")).strip()
        return [text] if text else []

    lines: list[str] = []
    for paragraph in text_frame.paragraphs:
        text = clean_text("".join(run.text for run in paragraph.runs) or paragraph.text or "").strip()
        if not text:
            continue
        level = int(getattr(paragraph, "level", 0) or 0)
        if level > 0:
            indent = "  " * min(level, 4)
            lines.append(f"{indent}- {text}")
        else:
            lines.append(text)
    return lines


class PptxParser(BaseParser):
    """PPTX 解析器。

    目标不是复刻完整页面样式，而是把 slide 里的标题和正文组织成可检索文本。
    """

    def parse(self, path: Path, source: str) -> Document:
        """把 PPTX 文档转成结构化文本。

        这里把 `pptx` 的导入放到函数内部，而不是模块顶层，
        是为了让项目在未安装 `python-pptx` 时仍能正常启动其它链路。

        只有真正入库 `.pptx` 文件时，才会显式要求这个依赖。
        """

        try:
            from pptx import Presentation
        except ModuleNotFoundError as exc:  # pragma: no cover - 依赖缺失属于环境问题，不是逻辑分支
            raise RuntimeError(
                "PPTX parsing requires `python-pptx`. Install it before ingesting .pptx files."
            ) from exc

        presentation = Presentation(str(path))
        lines: list[str] = []
        # 把第一页遇到的有效标题作为文档默认标题。
        # 这样在 PPT 文件名不规范时，返回给用户的标题也更自然。
        first_title = ""
        for slide_idx, slide in enumerate(presentation.slides, start=1):
            slide_title = _safe_slide_title(slide)
            if slide_title and not first_title:
                first_title = slide_title

            # 每个 slide 都显式写出一个标题块，便于后续切块时保留页内结构。
            if slide_title:
                lines.append(f"# Slide {slide_idx}: {slide_title}")
            else:
                lines.append(f"# Slide {slide_idx}")

            for shape in slide.shapes:
                # 标题已经单独处理过，避免重复写入。
                if getattr(slide.shapes, "title", None) is shape:
                    continue
                lines.extend(_shape_lines(shape))

        content = clean_text("\n".join(lines))
        return Document(
            # 统一由上游生成真正的文档唯一标识。
            doc_id="",
            source=source,
            title=first_title or path.stem,
            content=content,
            mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            # slide 总数对审计、展示和后续调试都很有用，先作为轻量 metadata 保留下来。
            metadata={"slides": len(presentation.slides)},
        )
