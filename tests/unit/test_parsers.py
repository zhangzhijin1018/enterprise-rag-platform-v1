"""文档解析器单元测试。

这一组测试关注第四轮增强新增的三类文档：

- `.txt`
- `.csv`
- `.pptx`

测试目标不是覆盖第三方库的全部行为，而是验证我们自己的解析契约：

1. 不同格式都能被统一成 `Document`
2. 结构化格式会按预期转成更利于检索的文本
3. 缺失可选依赖时，错误信息足够明确
4. `registry.py` 已经把新格式真正接入入库主链路
"""

from __future__ import annotations

import builtins
import importlib.util
from pathlib import Path
from typing import Any

import pytest

from core.ingestion.parsers.csv_parser import CsvParser
from core.ingestion.parsers.docx_parser import DocxParser
from core.ingestion.parsers.pptx_parser import PptxParser
from core.ingestion.parsers.registry import get_parser_for_filename
from core.ingestion.parsers.text_parser import TextParser


def test_text_parser_reads_plain_text_file(tmp_path: Path) -> None:
    """TXT 解析器应当直接保留纯文本可检索内容。"""

    file_path = tmp_path / "faq.txt"
    file_path.write_text("  第一行说明  \n\n第二行内容\t\t", encoding="utf-8")

    parser = TextParser()
    document = parser.parse(file_path, source="unit-test")

    assert document.title == "faq"
    assert document.mime_type == "text/plain"
    assert document.content == "第一行说明\n第二行内容"


def test_csv_parser_transforms_rows_into_structured_text(tmp_path: Path) -> None:
    """CSV 解析器应把表格行转成更适合 RAG 的字段化文本。"""

    file_path = tmp_path / "errors.csv"
    file_path.write_text(
        "code,reason,solution\n"
        "E-1001,Redis connection failed,Check network\n"
        "E-1002,Milvus timeout,Increase query timeout\n",
        encoding="utf-8",
    )

    parser = CsvParser()
    document = parser.parse(file_path, source="unit-test")

    assert document.mime_type == "text/csv"
    assert document.metadata["rows"] == 2
    assert document.metadata["columns"] == 3
    assert "## Row 1: E-1001" in document.content
    assert "code: E-1001" in document.content
    assert "solution: Increase query timeout" in document.content


@pytest.mark.skipif(
    importlib.util.find_spec("docx") is None,
    reason="python-docx 未安装，跳过真实 DOCX 解析测试",
)
def test_docx_parser_preserves_heading_levels_and_tables(tmp_path: Path) -> None:
    """DOCX 解析器应尽量保留标题层级、列表和表格结构。"""

    from docx import Document as WordDocument

    file_path = tmp_path / "manual.docx"
    document = WordDocument()
    document.add_heading("巡检制度", level=1)
    document.add_heading("巡检步骤", level=2)
    document.add_paragraph("先检查输煤皮带。", style="List Bullet")
    table = document.add_table(rows=2, cols=2)
    table.rows[0].cells[0].text = "code"
    table.rows[0].cells[1].text = "reason"
    table.rows[1].cells[0].text = "E-1001"
    table.rows[1].cells[1].text = "网络异常"
    document.save(file_path)

    parser = DocxParser()
    parsed = parser.parse(file_path, source="unit-test")

    assert "# 巡检制度" in parsed.content
    assert "## 巡检步骤" in parsed.content
    assert "- 先检查输煤皮带。" in parsed.content
    assert "## Table 1" in parsed.content
    assert "code: E-1001" in parsed.content


def test_registry_supports_txt_csv_and_pptx() -> None:
    """解析器注册表必须把新增格式真正暴露给入库主链路。"""

    assert isinstance(get_parser_for_filename("notes.txt"), TextParser)
    assert isinstance(get_parser_for_filename("faq.csv"), CsvParser)
    assert isinstance(get_parser_for_filename("training.pptx"), PptxParser)


def test_pptx_parser_raises_clear_error_when_dependency_missing(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    """未安装 `python-pptx` 时，应返回清晰可执行的错误，而不是栈追踪噪声。"""

    file_path = tmp_path / "deck.pptx"
    file_path.write_bytes(b"not-a-real-pptx")

    original_import = builtins.__import__

    def fake_import(
        name: str,
        globals: dict[str, Any] | None = None,
        locals: dict[str, Any] | None = None,
        fromlist: tuple[str, ...] = (),
        level: int = 0,
    ) -> Any:
        if name == "pptx":
            raise ModuleNotFoundError("No module named 'pptx'")
        return original_import(name, globals, locals, fromlist, level)

    monkeypatch.setattr(builtins, "__import__", fake_import)

    parser = PptxParser()
    with pytest.raises(RuntimeError, match="python-pptx"):
        parser.parse(file_path, source="unit-test")


@pytest.mark.skipif(
    importlib.util.find_spec("pptx") is None,
    reason="python-pptx 未安装，跳过真实 PPTX 解析测试",
)
def test_pptx_parser_extracts_slide_title_and_body(tmp_path: Path) -> None:
    """安装 `python-pptx` 时，验证 slide 标题和正文都能被提取。"""

    from pptx import Presentation

    file_path = tmp_path / "deck.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "RAG 方案介绍"
    text_frame = slide.placeholders[1].text_frame
    text_frame.text = "支持多路查询"
    paragraph = text_frame.add_paragraph()
    paragraph.level = 1
    paragraph.text = "支持父子切块"
    presentation.save(file_path)

    parser = PptxParser()
    document = parser.parse(file_path, source="unit-test")

    assert document.title == "RAG 方案介绍"
    assert document.metadata["slides"] == 1
    assert "# Slide 1: RAG 方案介绍" in document.content
    assert "支持多路查询" in document.content
    assert "- 支持父子切块" in document.content
